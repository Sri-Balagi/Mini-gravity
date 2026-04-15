import os
import subprocess
import re
import shutil
import requests

BASE_DIR=os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUTPUT_DIR=os.path.join(BASE_DIR,"output")
INPUT_DIR=os.path.join(BASE_DIR,"input")
file_path = os.path.join(BASE_DIR,"input","notes.pdf")


def extract_filename(details:dict):
    filename=details.get("filename") or details.get("file_name") or details.get("file name") or details.get("name")
    return filename

def extract_instruction(details:dict):
    instruction = details.get("instruction") or details.get("task")
    return instruction

def clean_code(raw: str) -> str:
    """Extract only code from a model response that may include markdown fences and explanatory text."""
    # Strip ANSI escape codes
    raw = re.sub(r'\x1b\[[0-9;]*[a-zA-Z]', '', raw)

    # Case 1: model wrapped code in a markdown fence  (```python ... ```)
    # Use DOTALL so `.` matches newlines too
    fence_match = re.search(r'```(?:python|py|sh|bash|)?\s*\r?\n(.*?)```', raw, re.DOTALL | re.IGNORECASE)
    if fence_match:
        return fence_match.group(1).strip()

    # Case 2: no fence — remove any stray ``` markers and return the rest
    lines = raw.splitlines()
    cleaned = [l for l in lines if not re.match(r'^\s*```', l)]
    return '\n'.join(cleaned).strip()


def valid_code(code:str):
    try:
        compile(code,"<string>","exec")
        return True
    except:
        return False
def create_file(details:dict):
    filename=extract_filename(details)
    instruction=extract_instruction(details) or ""

    if not filename or filename.strip()=="":
        raise Exception("File name not found or invalid file name")
    
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    filename=os.path.basename(filename)
    file_path=os.path.join(OUTPUT_DIR,filename)
    
    if "folder" in instruction.lower() or "directory" in instruction.lower():
        os.makedirs(file_path, exist_ok=True)
        return f"Folder created {file_path}"
    else:
        with open(file_path,"w") as f:
            pass
        return f"File created {file_path}"


def delete_file(details):
    file_path = resolve_path(details)
    if not file_path or not os.path.exists(file_path):
        raise Exception("File not found")
    if os.path.isdir(file_path):
        shutil.rmtree(file_path)
    else:
        os.remove(file_path)
    return f"Deleted {file_path}"
    
def rename_file(details):
    old_path = resolve_path(details, "old_name")
    new_name = os.path.basename(details.get("new_name", "renamed_file"))
    new_path = os.path.join(OUTPUT_DIR, new_name)

    if not old_path or not os.path.exists(old_path):
        raise Exception("Original File Not Found")
    os.rename(old_path, new_path)
    return f"Renamed {old_path} to {new_path}"

def move_file(details):
    filename=extract_filename(details)
    destination=details.get("destination", "")
    
    filename=os.path.basename(filename)
    dest_name=os.path.basename(destination)
    
    src_path=os.path.join(OUTPUT_DIR,filename)
    dest_path=os.path.join(OUTPUT_DIR,dest_name)
    
    if not os.path.exists(src_path):
        raise Exception("Source file not found")
        
    shutil.move(src_path, dest_path)
    return f"Moved {src_path} to {dest_path}"

def modify_file(details):
    instruction = extract_instruction(details)
    file_path = resolve_path(details)
    
    if not file_path or not os.path.exists(file_path):
        raise Exception("File to modify not found")
    
    with open(file_path,"r",encoding="utf-8") as f:
        content=f.read()

    prompt=f"""Modify the following code based on instruction.
    INSTRUCTION:
    {instruction}

    code:
    {content}
    """

    try:
        response = requests.post(
            "http://localhost:11434/api/generate",
            json={"model": "deepseek-coder:6.7B", "prompt": prompt, "stream": False},
            timeout=300
        )
        response.raise_for_status()
        updated_code = clean_code(response.json().get("response", ""))
    except Exception as e:
        raise Exception(f"Ollama API Error: {e}")

    with open(file_path,"w",encoding="utf-8") as f:
        f.write(updated_code)
    return f"updated {file_path}"

def resolve_path(details, filename_key="filename"):
    injected = details.get("file_path")
    if injected and os.path.exists(injected):
        return injected
    
    filename = details.get(filename_key) or extract_filename(details)
    if filename:
        return os.path.join(OUTPUT_DIR, os.path.basename(filename))
    return None

def _derive_filename(instruction: str) -> str:
    """Derive a safe snake_case .py filename from the instruction text."""
    # Take the first 5 meaningful words, lowercase, join with underscores
    words = re.sub(r'[^a-zA-Z0-9\s]', '', instruction).lower().split()
    stop = {"a","an","the","me","please","using","with","in","to","for","i","you","that","is","want"}
    keywords = [w for w in words if w not in stop][:4]
    name = "_".join(keywords) if keywords else "program"
    return name + ".py"

def write_code(details:dict,retries=3):
    filename=extract_filename(details)
    instruction=extract_instruction(details)
    
    if not instruction:
        raise Exception("Instruction not found — please describe what to build.")

    # Auto-generate filename from instruction if LLM didn't supply one
    if not filename or not filename.strip():
        filename = _derive_filename(instruction)
        print(f"No filename provided by model; derived: {filename}")

    #file creation
    filename=os.path.basename(filename)
    os.makedirs(OUTPUT_DIR,exist_ok=True)
    file_path=os.path.join(OUTPUT_DIR,filename)
    
    #prompt for making the LLM write code
    gui_directive = ""

    prompt = f"""You are a senior software engineer. Write complete, production-quality code.

USER REQUEST:
{instruction}

RULES — READ CAREFULLY:
1. Analyse the request fully before writing a single line.
2. Choose the right approach, language, and framework based on the user's request.
3. The code must be 100% complete and immediately runnable — no placeholders,
   no "# TODO", no truncated logic.
4. Every imported module must be used.  Every function must have a body.
5. Handle obvious runtime errors (division by zero, empty input, etc.).
6. Use clear, descriptive variable and function names.
7. DO NOT hardcode test values! If the program requires values to compute, you MUST accept them interactively using the `input()` function so the user can test different parameters.

OUTPUT FORMAT — CRITICAL:
- Return ONLY raw source code.
- Do NOT include markdown fences (```), explanations, or comments outside the code.
- The very first character of your response must be the first character of the code.
"""

    try:
        response = requests.post(
            "http://localhost:11434/api/generate",
            json={"model": "deepseek-coder:6.7B", "prompt": prompt, "stream": False},
            timeout=300
        )
        response.raise_for_status()
        raw_stdout = response.json().get("response", "").strip()
    except Exception as e:
        raise Exception(f"Ollama API Error: {e}")

    code = clean_code(raw_stdout)

    # checking for code correctness by execution check
    if not valid_code(code):
        print("Attempting repair")
        code = repair_code(code)

    if not code.strip():
        code = "# Code generation failed. Please try again with a more specific prompt."

    #writing to file
    with open(file_path,"w",encoding="utf-8") as f:
        f.write(code)
    return f"code written to {file_path}" 

def repair_code(code):
    prompt=f"""Fix the following code. It contains syntax or formatting errors.

    Return ONLY corrected code and nothing else.

    CODE:
    {code}
    """
    try:
        response = requests.post(
            "http://localhost:11434/api/generate",
            json={"model": "deepseek-coder:6.7B", "prompt": prompt, "stream": False},
            timeout=300
        )
        response.raise_for_status()
        return clean_code(response.json().get("response", ""))
    except Exception as e:
        raise Exception(f"Ollama API Error: {e}")

def repair_file(details:dict):
    file_path = resolve_path(details)
    if not file_path or not os.path.exists(file_path):
        raise Exception("File to repair not found")
    with open(file_path,"r",encoding="utf-8") as f:
        code=f.read()
    fixed_code=repair_code(code)
    with open(file_path,"w",encoding="utf-8") as f:
        f.write(fixed_code)
    return f"Repaired file:{file_path}"


def summarize_text(details,filename):
    pass

def extract_txt(path):
    with open (path,"r", encoding="utf-8") as f:
        return f.read()

def extract_pdf(path):
    import PyPDF2
    text=""
    with open(path,"rb") as f:
        reader=PyPDF2.PdfReader(f)
        for page in reader.pages:
            text+=page.extract_text() or ""
    return text

def extract_docx(path):
    from docx import Document
    doc=Document(path)
    return "\n".join([p.text for p in doc.paragraphs])

def extract_pptx(path):
    from pptx import Presentation
    prs=Presentation(path)
    text=""
    for slide in prs.slides:
        for shape in slide.shapes:
                if hasattr(shape,"text"):
                    text+=shape.text + "\n"
    return text

def extract_xlsx(path):
    import pandas as pd
    df=pd.read_excel(path)
    return df.to_string()

def extract_text_from_file(file_path):
    ext=file_path.lower().split(".")[-1]
    if ext=="txt":
        return extract_txt(file_path)
    elif ext=="pdf":
        return extract_pdf(file_path)
    elif ext=="docx":
        return extract_docx(file_path)
    elif ext=="pptx":
        return extract_pptx(file_path)
    elif ext=="xlsx":
        return extract_xlsx(file_path)
    else:
        return "Unsopperted file type: {ext}"
    
def summarize(details):
    filename = details.get("filename")
    injected_path = details.get("file_path")   # absolute path injected by the UI
    text = None
    filepath = None

    # Priority 1: absolute file_path injected by the UI (uploaded document)
    if injected_path and os.path.exists(injected_path):
        text = extract_text_from_file(injected_path)
        filepath = injected_path

    # Priority 2: filename provided by the LLM → resolve against output dir
    elif filename and filename.strip():
        filepath = os.path.join(OUTPUT_DIR, os.path.basename(filename))
        if os.path.exists(filepath):
            text = extract_text_from_file(filepath)
        else:
            raise Exception(f"File not found: {filepath}")

    # Priority 3: raw text / instruction passed directly
    if not text:
        text = details.get("text") or details.get("instruction")

    if not text:
        raise Exception("No content to summarize — please upload a file or provide text.")

    prompt = f"""Summarize the following text accurately.
Highlight key sections and important details only.
{text}"""

    try:
        response = requests.post(
            "http://localhost:11434/api/generate",
            json={"model": "qwen2.5:7b", "prompt": prompt, "stream": False},
            timeout=300
        )
        response.raise_for_status()
        return response.json().get("response", "").strip()
    except Exception as e:
        raise Exception(f"Ollama API Error: {e}")

def general_chat(details):
    user_input = details.get("text") or details.get("query") or details.get("instruction") or details.get("user_message")
    if not user_input and details:
        user_input = " ".join(str(v) for v in details.values())
    if not user_input:
        raise Exception("No input for chat")
    try:
        response = requests.post(
            "http://localhost:11434/api/generate",
            json={"model": "qwen2.5:7b", "prompt": user_input, "stream": False},
            timeout=300
        )
        response.raise_for_status()
        return response.json().get("response", "").strip()
    except Exception as e:
        raise Exception(f"Ollama API Error: {e}")

def execute(intent_data:dict):
    intent=intent_data.get("intent")
    details=intent_data.get("details",{})

    if intent=="create_file":
        return create_file(details)
    elif intent=="write_code":
        return write_code(details)
    elif intent=="delete_file":
        return delete_file(details)
    elif intent=="rename_file":
        return rename_file(details)
    elif intent=="move_file":
        return move_file(details)
    elif intent=="modify_file":
        return modify_file(details)
    elif intent=="repair_file":
        return repair_file(details)
    elif intent=="summarize":
        return summarize(details)
    elif intent=="general_chat":
        return general_chat(details)
    else:
        raise Exception(f"Unknown intent:{intent}")

